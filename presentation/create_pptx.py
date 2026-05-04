#!/usr/bin/env python3
"""Generate MediGraph PowerPoint presentation - detailed version matching report quality."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x16, 0x21, 0x3e)
ACCENT_BLUE = RGBColor(0x0f, 0x4c, 0x75)
LIGHT_BLUE = RGBColor(0x3b, 0x82, 0xf6)
RED = RGBColor(0xdc, 0x35, 0x45)
ORANGE = RGBColor(0xff, 0xc1, 0x07)
GREEN = RGBColor(0x28, 0xa7, 0x45)
GRAY = RGBColor(0x4a, 0x5a, 0x6a)
LIGHT_GRAY = RGBColor(0xf8, 0xf9, 0xfa)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x1a, 0x1a, 0x2e)


def add_title_slide(prs, title, subtitle, footer):
    """Add a title slide with gradient-style design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.733), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Footer with names
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.733), Inches(1.2))
    tf = footer_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = footer
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_slide_header(slide, title, prs):
    """Add consistent header to content slides."""
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12.133), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_metric_boxes(slide, metrics, y_pos):
    """Add metric boxes row."""
    box_width = 2.4
    total_width = len(metrics) * box_width + (len(metrics) - 1) * 0.3
    start_x = (13.333 - total_width) / 2

    for i, (value, label) in enumerate(metrics):
        x = Inches(start_x + i * (box_width + 0.3))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y_pos), Inches(box_width), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_BLUE
        box.line.fill.background()

        # Value
        val_box = slide.shapes.add_textbox(x, Inches(y_pos + 0.15), Inches(box_width), Inches(0.55))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x, Inches(y_pos + 0.65), Inches(box_width), Inches(0.35))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xcc, 0xdd, 0xee)
        p.alignment = PP_ALIGN.CENTER


def add_table(slide, headers, rows, x, y, width, row_height=0.38):
    """Add a formatted table."""
    cols = len(headers)
    row_count = len(rows) + 1
    col_width = width / cols

    table = slide.shapes.add_table(row_count, cols, Inches(x), Inches(y), Inches(width), Inches(row_height * row_count)).table

    # Set column widths
    for i in range(cols):
        table.columns[i].width = Inches(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(11)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].space_before = Pt(3)
        cell.text_frame.paragraphs[0].space_after = Pt(3)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = GRAY
            para.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            para.space_before = Pt(2)
            para.space_after = Pt(2)
            # Alternate row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table


def add_text_box(slide, text, x, y, width, height, font_size=12, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box with formatting."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box


def add_bullet_list(slide, items, x, y, width, height, font_size=12, color=GRAY):
    """Add a bullet list."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return box


# ============ SLIDE 1: Title ============
add_title_slide(
    prs,
    "Provera",
    "AI-Powered Medicare Fraud Detection Using Hybrid ML + Agentic AI",
    "Nikita Ravi  |  Jagannath Narayanswamy\nData Analytics Capstone | Spring 2026 | George Washington University"
)

# ============ SLIDE 2: Executive Summary ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Executive Summary", prs)

# Key metrics
add_metric_boxes(slide, [
    ("0.91", "ROC-AUC"),
    ("7/7", "Golden Set"),
    ("11,090", "Facilities"),
    ("85%", "Precision@100"),
], 1.4)

# Summary text
add_text_box(slide,
    "Provera combines graph-based machine learning with LLM agents to detect Medicare fraud rings. "
    "The system analyzes ownership structures, shared addresses, and billing patterns to identify "
    "coordinated fraud schemes that evade traditional rule-based detection.",
    0.6, 2.8, 12.133, 0.9, font_size=14)

# Two column layout
add_text_box(slide, "What We Detect", 0.6, 3.8, 5.5, 0.4, font_size=14, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, [
    "Shell company networks (entity cycling)",
    "Address/phone sharing schemes",
    "Connections to excluded providers",
    "Ownership concentration patterns"
], 0.6, 4.2, 5.5, 2)

add_text_box(slide, "Key Innovation", 6.8, 3.8, 5.5, 0.4, font_size=14, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, [
    "Structural detection, not just billing anomalies",
    "AI agent generates investigation briefs",
    "DOJ cross-reference catches behavioral fraud",
    "100% factual accuracy on generated reports"
], 6.8, 4.2, 5.5, 2)


# ============ SLIDE 3: Problem Statement ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "The Problem: $60-100 Billion in Annual Medicare Fraud", prs)

# Stats boxes
stats = [
    ("$60-100B", "Annual fraud losses"),
    ("$17B", "HHA improper payments (2023)"),
    ("40%", "FL share of prosecutions"),
    ("90%+", "False positive rate (current)"),
]
add_metric_boxes(slide, stats, 1.4)

# Problem details
add_text_box(slide, "Why Current Detection Fails", 0.6, 2.8, 5.8, 0.4, font_size=14, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, [
    "Rule-based systems trigger on billing anomalies AFTER fraud occurs",
    "Fraud rings operate through multiple shell companies that appear independent",
    "No visibility into ownership networks or shared infrastructure",
    "By the time alerts trigger, millions have already been stolen"
], 0.6, 3.2, 5.8, 2.5, font_size=12)

add_text_box(slide, "Our Approach: Structural Detection", 6.8, 2.8, 5.8, 0.4, font_size=14, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, [
    "Model provider networks as graphs with ownership edges",
    "Detect fraud rings BEFORE billing anomalies emerge",
    "Use graph ML to find hidden connections between entities",
    "Deploy AI agents to investigate and generate human-readable reports"
], 6.8, 3.2, 5.8, 2.5, font_size=12)

# Florida callout
callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.5), Inches(12.133), Inches(1.2))
callout.fill.solid()
callout.fill.fore_color.rgb = RGBColor(0xff, 0xf3, 0xcd)
callout.line.color.rgb = ORANGE

add_text_box(slide, "Florida Focus: Despite having only 7% of Medicare beneficiaries, Florida accounts for 40% of all "
    "Medicare fraud prosecutions. Home Health Agencies are the #1 fraud vector, making Florida HHAs "
    "the ideal testbed for fraud detection systems.", 0.8, 5.65, 11.733, 0.9, font_size=12, color=RGBColor(0x85, 0x64, 0x04))


# ============ SLIDE 4: Architecture Evolution ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Architecture Evolution: From 0.72 to 0.91 ROC-AUC", prs)

add_text_box(slide, "Each iteration addressed limitations discovered in the previous version:",
    0.6, 1.3, 12.133, 0.4, font_size=13, color=GRAY)

add_table(slide,
    ["Iteration", "Components Added", "ROC-AUC", "Problem Solved"],
    [
        ("1. Baseline ML", "XGBoost on billing features only", "0.72", "—"),
        ("2. + Graph", "Network centrality, community detection", "0.81", "Missing network context"),
        ("3. + Ownership", "Shared owners, addresses, phones", "0.87", "Invisible ownership patterns"),
        ("4. + AI Agent", "LLM investigation, DOJ cross-ref", "0.91*", "No explainability"),
    ],
    0.6, 1.9, 12.133, row_height=0.5
)

add_text_box(slide, "*Agent layer adds 100% factual accuracy on generated investigation reports",
    0.6, 4.2, 12.133, 0.3, font_size=10, color=RGBColor(0x88, 0x88, 0x88))

# Key insight box
insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.7), Inches(12.133), Inches(2.2))
insight.fill.solid()
insight.fill.fore_color.rgb = RGBColor(0xf0, 0xf7, 0xff)
insight.line.color.rgb = LIGHT_BLUE

add_text_box(slide, "Key Insight: Why Hybrid?", 0.8, 4.9, 11.733, 0.35, font_size=14, bold=True, color=ACCENT_BLUE)
add_text_box(slide,
    "ML alone: Fast and scalable, but high false positives and no context. Flags legitimate hospital systems.\n"
    "Agent alone: Explainable and contextual, but slow, expensive, and inconsistent at scale.\n"
    "Hybrid: ML triages 11,090 facilities instantly; Agent investigates only high-risk cases (~$0.03 each).",
    0.8, 5.3, 11.733, 1.4, font_size=12, color=GRAY)


# ============ SLIDE 5: System Architecture ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "System Architecture", prs)

# Architecture layers
layers = [
    ("DATA LAYER", "NPPES (11K providers) + LEIE (298 excluded) + Medicare Billing + FL SunBiz", RGBColor(0x64, 0x74, 0x8b)),
    ("GRAPH LAYER", "Provider Network: 26K nodes, 38K edges | Louvain: 8,015 communities", RGBColor(0x8b, 0x5c, 0xf6)),
    ("ML LAYER", "XGBoost (30+ features) | ROC-AUC: 0.91 | SHAP Explainability", RGBColor(0x3b, 0x82, 0xf6)),
    ("AGENT LAYER", "Claude 3.5 Sonnet | Red Flags (5 checks) | DOJ Cross-Ref | Hypothesis Gen/Eval", RGBColor(0x22, 0xc5, 0x5e)),
    ("OUTPUT", "Classification: HIGH/MEDIUM/LOW/CLEARED | Investigation Dossier | Recommendations", RGBColor(0xf5, 0x9e, 0x0b)),
]

y_start = 1.4
layer_height = 0.95
gap = 0.12

for i, (name, desc, color) in enumerate(layers):
    y = Inches(y_start + i * (layer_height + gap))

    # Layer box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.133), Inches(layer_height))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    # Layer name
    name_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.05), Inches(2.2), Inches(layer_height - 0.1))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    name_box.text_frame.anchor = MSO_ANCHOR.MIDDLE

    # Separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.1), y + Inches(0.15), Inches(0.02), Inches(layer_height - 0.3))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    line.line.fill.background()

    # Description
    desc_box = slide.shapes.add_textbox(Inches(3.3), y + Inches(0.05), Inches(9.2), Inches(layer_height - 0.1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    desc_box.text_frame.anchor = MSO_ANCHOR.MIDDLE

    # Arrow (except last)
    if i < len(layers) - 1:
        arrow_y = y + Inches(layer_height)
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), arrow_y, Inches(0.5), Inches(gap))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        arrow.line.fill.background()


# ============ SLIDE 6: Data Sources ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Data Sources: 5 Integrated Datasets", prs)

# Data sources table with details
add_table(slide,
    ["Source", "Records", "Key Fields", "Purpose"],
    [
        ("CMS NPPES", "11,090 HHAs", "NPI, name, address, phone", "Provider identity"),
        ("OIG LEIE", "298 excluded", "Exclusion type, date, NPI", "Ground truth labels"),
        ("Medicare Billing", "11,090", "Charges, payments, beneficiaries", "Billing behavior"),
        ("FL SunBiz", "Lookups", "Incorporation date, principals", "Entity age + ownership"),
        ("Provider Enrollment", "11,090", "Ownership %, role, dates", "Ownership network"),
    ],
    0.6, 1.4, 12.133, row_height=0.48
)

# Key fields breakdown
add_text_box(slide, "Critical Fields for Fraud Detection", 0.6, 4.0, 12.133, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)

datasets = [
    ("NPPES", "NPI (unique ID), Address (co-location), Phone (shared ops), Enumeration Date (entity age)"),
    ("LEIE", "Exclusion Type (fraud category), Exclusion Date (timing), Reinstatement (still active?)"),
    ("SunBiz", "Date Filed (shell company detection - entities <3 years old are high risk)"),
    ("Enrollment", "Owner Name (network links), Ownership % (control), Multiple NPIs (concentration)"),
]

y = 4.4
for name, fields in datasets:
    add_text_box(slide, name + ":", 0.6, y, 1.4, 0.35, font_size=11, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, fields, 2.1, y, 10.5, 0.35, font_size=10, color=GRAY)
    y += 0.4

# Data quality note
note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.0), Inches(12.133), Inches(0.8))
note.fill.solid()
note.fill.fore_color.rgb = RGBColor(0xe3, 0xf2, 0xfd)
note.line.color.rgb = LIGHT_BLUE

add_text_box(slide, "Data Quality: All datasets are from official government sources (CMS, OIG, FL Division of Corporations). "
    "Entity resolution performed via NPI matching (exact) and fuzzy name+address matching (Levenshtein distance <3).",
    0.8, 6.15, 11.733, 0.55, font_size=10, color=ACCENT_BLUE)


# ============ SLIDE 7: Graph Data Model ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Graph Data Model: Entities & Relationships", prs)

# Entity types (left side)
add_text_box(slide, "Entity Types (Nodes)", 0.6, 1.35, 5.5, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)

entities = [
    ("FACILITY", "11,090", "NPI, name, address, risk_score, is_excluded", LIGHT_BLUE),
    ("OWNER", "5,854", "owner_id, name, type (individual/org)", RGBColor(0x8b, 0x5c, 0xf6)),
    ("ADDRESS", "~9,800", "address_hash, street, city, zip", GREEN),
    ("PHONE", "~9,796", "phone_number (normalized)", ORANGE),
]

y = 1.8
for name, count, fields, color in entities:
    # Entity box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(5.5), Inches(0.65))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    add_text_box(slide, f"{name} ({count})", 0.75, y + 0.08, 2.5, 0.25, font_size=11, bold=True, color=WHITE)
    add_text_box(slide, fields, 0.75, y + 0.35, 5.2, 0.25, font_size=9, color=WHITE)
    y += 0.75

# Relationship types (right side)
add_text_box(slide, "Relationship Types (Edges)", 6.5, 1.35, 6.0, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)

relationships = [
    ("OWNS", "27,002", "OWNER → FACILITY", "ownership_%, role, start_date"),
    ("LOCATED_AT", "11,090", "FACILITY → ADDRESS", "primary location"),
    ("HAS_PHONE", "11,090", "FACILITY → PHONE", "contact number"),
    ("SHARES_ADDRESS", "3,170", "FACILITY ↔ FACILITY", "derived: same address"),
    ("SHARES_OWNER", "35,122", "FACILITY ↔ FACILITY", "derived: common owner"),
]

y = 1.8
for name, count, direction, attrs in relationships:
    add_text_box(slide, f"{name}", 6.5, y, 2.2, 0.25, font_size=10, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, f"({count})", 8.7, y, 0.8, 0.25, font_size=9, color=GRAY)
    add_text_box(slide, direction, 9.5, y, 2.0, 0.25, font_size=9, color=GRAY)
    add_text_box(slide, attrs, 6.5, y + 0.25, 6.0, 0.25, font_size=8, color=RGBColor(0x88, 0x88, 0x88))
    y += 0.55

# Graph statistics
add_text_box(slide, "Graph Statistics", 0.6, 5.0, 12.133, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)

add_table(slide,
    ["Metric", "Value", "Metric", "Value"],
    [
        ("Total Nodes", "26,002", "Louvain Communities", "8,015"),
        ("Total Edges", "37,691", "Avg Clustering Coef", "0.42"),
        ("Facility Nodes", "11,090", "Network Density", "0.0007"),
    ],
    0.6, 5.4, 12.133, row_height=0.38
)

# Community detection note
add_text_box(slide, "Community Detection: Louvain algorithm identifies natural clusters. Fraud rings operate as communities - "
    "investigating the network (not just individuals) reveals hidden connections.",
    0.6, 6.65, 12.133, 0.35, font_size=10, color=RGBColor(0x66, 0x77, 0x88))


# ============ SLIDE 8: ML Model & Features ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "ML Model & Feature Engineering", prs)

# Model comparison table
add_text_box(slide, "Model Selection", 0.6, 1.3, 5.8, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)
add_table(slide,
    ["Model", "ROC-AUC", "Prec@100"],
    [
        ("Random Forest", "0.84", "72%"),
        ("XGBoost", "0.91", "85%"),
        ("LightGBM", "0.89", "81%"),
        ("GCN (2-layer)", "0.88", "79%"),
    ],
    0.6, 1.7, 4.0, row_height=0.4
)

# Highlight XGBoost row
add_text_box(slide, "XGBoost selected for best accuracy + SHAP support",
    0.6, 3.55, 4.0, 0.3, font_size=9, color=RGBColor(0x22, 0xc5, 0x5e))

# Feature categories
add_text_box(slide, "30+ Features Across 4 Categories", 5.0, 1.3, 7.733, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)

# Features in columns
add_text_box(slide, "Billing (6)", 5.0, 1.75, 3.5, 0.3, font_size=11, bold=True, color=GRAY)
add_bullet_list(slide, ["Charges, payments", "Beneficiary count", "Charges per beneficiary", "Payment ratio"],
    5.0, 2.05, 3.5, 1.5, font_size=10)

add_text_box(slide, "Network (8)", 8.6, 1.75, 4.0, 0.3, font_size=11, bold=True, color=GRAY)
add_bullet_list(slide, ["PageRank, degree centrality", "Community size", "Fraud density in community", "Excluded neighbor ratio"],
    8.6, 2.05, 4.0, 1.5, font_size=10)

add_text_box(slide, "Ownership (8)", 5.0, 3.65, 3.5, 0.3, font_size=11, bold=True, color=GRAY)
add_bullet_list(slide, ["Owner count", "Max facilities/owner", "Shared address count", "Entity age (years)"],
    5.0, 3.95, 3.5, 1.5, font_size=10)

add_text_box(slide, "Derived (8+)", 8.6, 3.65, 4.0, 0.3, font_size=11, bold=True, color=GRAY)
add_bullet_list(slide, ["Community excluded count", "Fraud neighbor ratio", "Billing deviation", "Ownership concentration"],
    8.6, 3.95, 4.0, 1.5, font_size=10)

# Top SHAP features
add_text_box(slide, "Top SHAP Features (Risk Drivers)", 0.6, 4.0, 4.2, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)
add_table(slide,
    ["Feature", "Impact", "Direction"],
    [
        ("shared_address_count", "+0.32", "Higher = Risk"),
        ("community_excluded_count", "+0.28", "Higher = Risk"),
        ("entity_age_years", "+0.10", "Lower = Risk"),
    ],
    0.6, 4.4, 4.0, row_height=0.38
)


# ============ SLIDE 9: AI Agent Pipeline ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "AI Agent Investigation Pipeline", prs)

# Pipeline steps
steps = [
    ("1", "ML TRIAGE", "XGBoost scores all 11,090 facilities\nFilters to ~500 high-risk communities\nCost: ~$0 (local inference)", LIGHT_BLUE),
    ("2", "RED FLAGS", "5 deterministic checks per community\nNo LLM calls - pure rule-based\nOwnership, address, LEIE, billing, phone", RGBColor(0x8b, 0x5c, 0xf6)),
    ("3", "DOJ CHECK", "Cross-reference against prosecution DB\nForces full investigation if match\nCatches behavioral fraud ML misses", ORANGE),
    ("4", "AI INVESTIGATION", "Claude 3.5 Sonnet analysis\n3 LLM calls: hypotheses, evaluation, narrative\nCost: ~$0.03 per investigation", GREEN),
]

x_start = 0.5
box_width = 3.0
gap = 0.25

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(x_start + i * (box_width + gap))

    # Step box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(box_width), Inches(2.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = color

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), Inches(1.55), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    num_box = slide.shapes.add_textbox(x + Inches(1.2), Inches(1.6), Inches(0.55), Inches(0.45))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, title, x_start + i * (box_width + gap) + 0.15, 2.2, box_width - 0.3, 0.4,
        font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, desc, x_start + i * (box_width + gap) + 0.15, 2.6, box_width - 0.3, 1.4,
        font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Red flag checklist
add_text_box(slide, "5 Automated Red Flags (Deterministic)", 0.6, 4.5, 12.133, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)
add_table(slide,
    ["#", "Red Flag", "Trigger Condition"],
    [
        ("1", "Ownership Concentration", "Single owner controls 3+ facilities"),
        ("2", "Shared Address", "Multiple entities at same address"),
        ("3", "LEIE Connection", "Any facility/owner on exclusion list"),
        ("4", "Billing Deviation", "Billing > 2σ from state median"),
        ("5", "Shared Phone", "Multiple entities sharing phone"),
    ],
    0.6, 4.9, 12.133, row_height=0.38
)


# ============ SLIDE 10: Case Studies ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Case Studies: 5 Classification Types", prs)

cases = [
    ("HIGH", "Shell Company Ring (Community 1597)",
     "4 HHAs, 100% excluded, similar Miami addresses",
     "\"Entity cycling to evade oversight. All 4 received NPPES-linked exclusions.\"", RED),
    ("LOW", "Legitimate Hospital System (Community 170)",
     "227 facilities, Kindred Healthcare, 1.3% exclusion rate",
     "\"PE-backed chain with legitimate infrastructure. Shared structure is NOT fraud.\"", GREEN),
    ("CLEARED", "Established Nonprofit (Community 215)",
     "52-year-old organization, zero exclusions",
     "\"Long operational history rules out shell company classification.\"", GREEN),
    ("HIGH", "DOJ-Prosecuted (NPI 1851563381)",
     "Florida Patient Care Corp, phantom billing prosecution",
     "\"DOJ RECORD FOUND. Prosecuted for billing for services never provided.\"", RED),
    ("MEDIUM", "Doral Fraud Corridor (Community 731)",
     "5 HHAs co-located, 1/5 excluded, DOJ Strike Force area",
     "\"Geographic risk significant. Cannot confirm fraud without claims data.\"", ORANGE),
]

y = 1.35
for classification, title, evidence, finding, color in cases:
    # Classification badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(0.9), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.02), Inches(0.9), Inches(0.3))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = classification
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, title, 1.6, y, 4.0, 0.35, font_size=12, bold=True, color=BLACK)

    # Evidence
    add_text_box(slide, evidence, 5.7, y, 3.2, 0.35, font_size=10, color=GRAY)

    # Finding (italic quote)
    add_text_box(slide, finding, 9.0, y, 3.9, 0.35, font_size=9, color=RGBColor(0x55, 0x66, 0x77))

    y += 0.55

# Summary box
summary = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.15), Inches(12.133), Inches(0.7))
summary.fill.solid()
summary.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)
summary.line.color.rgb = GREEN

add_text_box(slide, "Key: Agent explains WHY each classification was made, not just the score. "
    "This enables investigators to trust or override the recommendation with full context.",
    0.8, 6.28, 11.733, 0.5, font_size=11, color=RGBColor(0x2e, 0x7d, 0x32))


# ============ SLIDE 11: Edge Cases & Smart Thresholding ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Edge Cases: Why ML Score ≠ Final Classification", prs)

# Threshold flow diagram
add_text_box(slide, "Smart Investigation Thresholding", 0.6, 1.3, 12.133, 0.35, font_size=14, bold=True, color=ACCENT_BLUE)

# Flow boxes
flow_steps = [
    ("ML SCORES\nALL 11,090", LIGHT_BLUE, "~0.1s"),
    ("Score > 0.3?\nOR Flags > 0?\nOR DOJ Match?", ORANGE, "Filter"),
    ("YES → Full\nInvestigation", GREEN, "~$0.03"),
    ("NO → Skip\n(Save Cost)", GRAY, "$0"),
]

for i, (text, color, label) in enumerate(flow_steps):
    x = Inches(0.8 + i * 3.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(2.6), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    txt = slide.shapes.add_textbox(x, Inches(1.8), Inches(2.6), Inches(0.9))
    tf = txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Label below
    add_text_box(slide, label, 0.8 + i * 3.1, 2.85, 2.6, 0.3, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Arrow
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.65), Inches(2.1), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        arrow.line.fill.background()

# Edge cases table
add_text_box(slide, "Edge Cases Handled Correctly", 0.6, 3.3, 12.133, 0.35, font_size=14, bold=True, color=ACCENT_BLUE)

add_table(slide,
    ["Scenario", "ML Score", "Classification", "Why?"],
    [
        ("Kindred Healthcare (227 facilities)", "0.81", "LOW", "Agent recognizes legitimate PE-backed system"),
        ("Community 698 (large network)", "0.80", "LOW", "High score from size, but 0% exclusion rate"),
        ("Florida Patient Care Corp", "0.42", "HIGH", "DOJ match forces investigation despite low score"),
        ("Doral office building (5 HHAs)", "0.72", "MEDIUM", "Co-location is legitimate, but geographic risk"),
        ("New entity, clean record", "0.55", "LOW", "No red flags despite being new"),
    ],
    0.6, 3.7, 12.133, row_height=0.42
)

# Key insight
insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.0), Inches(12.133), Inches(0.9))
insight.fill.solid()
insight.fill.fore_color.rgb = RGBColor(0xff, 0xf8, 0xe1)
insight.line.color.rgb = ORANGE

add_text_box(slide, "Key: ML score is a starting point, not the answer. The agent evaluates context that ML cannot see: "
    "entity age, corporate structure legitimacy, DOJ history, and geographic risk factors. This reduces false positives by 60%+.",
    0.8, 6.15, 11.733, 0.6, font_size=11, color=RGBColor(0x8d, 0x6e, 0x00))


# ============ SLIDE 12: Iterative Improvements ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Iterative Improvements: Making the Agent Smarter", prs)

# Timeline of improvements
improvements = [
    ("1", "SunBiz Integration", "Added FL corporate records for entity age. Shell companies are NEW by definition—52-year entities cannot be shell companies.", GREEN),
    ("2", "DOJ Cross-Reference", "Built prosecution database (50+ cases). Forces investigation on matches regardless of ML score. Catches behavioral fraud.", GREEN),
    ("3", "Red Flag Checklist", "5 deterministic checks BEFORE LLM. Agent cannot contradict hard facts (exclusion status, shared addresses).", GREEN),
    ("4", "False Positive Rules", "System prompt warns about Kindred, HCA, Solaris. Large healthcare systems have legitimate shared infrastructure.", GREEN),
    ("5", "SHAP Explainability", "Every ML prediction shows top risk factors. Investigators see WHY something was flagged, not just the score.", GREEN),
    ("6", "Community Detection", "Louvain algorithm finds natural clusters. Fraud rings operate in communities—investigate the network, not just individuals.", GREEN),
    ("7", "Confidence Calibration", "Agent must state HIGH/MEDIUM/LOW confidence with reasoning. Uncertainty is surfaced, not hidden.", GREEN),
]

y = 1.35
for num, title, desc, color in improvements:
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y), Inches(0.4), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_BLUE
    badge.line.fill.background()

    num_text = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.05), Inches(0.4), Inches(0.3))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, title, 1.1, y + 0.05, 2.4, 0.35, font_size=12, bold=True, color=BLACK)

    # Description
    add_text_box(slide, desc, 3.6, y, 9.0, 0.5, font_size=10, color=GRAY)

    y += 0.58

# Future improvements box
future = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.55), Inches(12.133), Inches(1.35))
future.fill.solid()
future.fill.fore_color.rgb = RGBColor(0xf3, 0xe5, 0xf5)
future.line.color.rgb = RGBColor(0x9c, 0x27, 0xb0)

add_text_box(slide, "Next Iterations (Future Work)", 0.8, 5.7, 11.733, 0.3, font_size=12, bold=True, color=RGBColor(0x7b, 0x1f, 0xa2))
add_bullet_list(slide, [
    "Claims-level analysis: Detect phantom billing, upcoding, kickback patterns",
    "Temporal tracking: Monitor entity cycling and ownership changes over time",
    "Beneficiary network: Identify patient-sharing patterns across fraud rings"
], 0.8, 6.0, 11.733, 0.85, font_size=10, color=RGBColor(0x6a, 0x1b, 0x9a))


# ============ SLIDE 13: Golden Set Evaluation ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Golden Set Evaluation: 7/7 Accuracy", prs)

add_text_box(slide, "Curated test cases covering all classification types and edge cases:",
    0.6, 1.3, 12.133, 0.35, font_size=13, color=GRAY)

add_table(slide,
    ["Community", "Type", "Expected", "Actual", "Key Evidence"],
    [
        ("1597", "Shell company ring", "HIGH", "HIGH", "4/4 excluded, entity cycling"),
        ("731", "Doral fraud corridor", "MEDIUM", "MEDIUM", "1/5 excluded, geographic risk"),
        ("215", "Established nonprofit", "LOW", "LOW", "52-year entity, no exclusions"),
        ("170", "PE healthcare chain", "LOW", "LOW", "Kindred system, legitimate"),
        ("4446", "Mixed risk cluster", "MEDIUM", "MEDIUM", "2/6 excluded, shared owner"),
        ("5806", "Clean single facility", "CLEARED", "CLEARED", "0 flags, isolated node"),
        ("5411", "New entity, clean", "LOW", "LOW", "Recent but no red flags"),
    ],
    0.6, 1.7, 12.133, row_height=0.45
)

# Result boxes
add_metric_boxes(slide, [
    ("7/7", "Classification Accuracy"),
    ("100%", "Factual Accuracy"),
    ("7/8", "DOJ Cases Identified"),
    ("35s", "Avg Investigation Time"),
], 5.3)

# Note
add_text_box(slide, "Factual accuracy verified: All facility names, NPIs, addresses, and exclusion statuses "
    "cross-referenced against source data. Zero hallucinations in generated investigation briefs.",
    0.6, 6.6, 12.133, 0.4, font_size=11, color=RGBColor(0x66, 0x77, 0x88))


# ============ SLIDE 14: Limitations & Scalability ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "Limitations & Scalability", prs)

# What we detect vs can't detect
add_text_box(slide, "What We Detect (Structural Fraud)", 0.6, 1.35, 5.8, 0.35, font_size=13, bold=True, color=GREEN)
add_bullet_list(slide, [
    "Shell company networks and entity cycling",
    "Address/phone sharing schemes",
    "Connections to excluded providers",
    "Ownership concentration patterns"
], 0.6, 1.75, 5.8, 1.8, font_size=12)

add_text_box(slide, "What We Cannot Detect (Behavioral)", 6.8, 1.35, 5.8, 0.35, font_size=13, bold=True, color=RED)
add_bullet_list(slide, [
    "Kickback schemes (requires referral data)",
    "Phantom billing (requires claims detail)",
    "Upcoding (requires procedure analysis)",
    "3/7 DOJ cases = behavioral fraud we'd miss"
], 6.8, 1.75, 5.8, 1.8, font_size=12)

# Scalability
add_text_box(slide, "Scalability: Florida to National", 0.6, 3.6, 12.133, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)
add_table(slide,
    ["Metric", "Current (Florida)", "National Scale"],
    [
        ("Providers", "11,090", "1.2M+"),
        ("Communities", "8,015", "~800K"),
        ("Investigation Cost", "$0.03 each", "~$36K full scan"),
    ],
    0.6, 4.0, 5.5, row_height=0.42
)

# Future enhancements
add_text_box(slide, "Future Enhancements", 6.8, 3.6, 5.8, 0.35, font_size=13, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, [
    "Claims integration for behavioral fraud",
    "Temporal entity cycling analysis",
    "Multi-state expansion (TX, CA, NY)",
    "Real-time new registration monitoring"
], 6.8, 4.0, 5.8, 2.0, font_size=12)

# Conclusion box
conclusion = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.7), Inches(12.133), Inches(1.2))
conclusion.fill.solid()
conclusion.fill.fore_color.rgb = RGBColor(0xf0, 0xf7, 0xff)
conclusion.line.color.rgb = ACCENT_BLUE

add_text_box(slide, "Conclusion", 0.8, 5.85, 11.733, 0.3, font_size=14, bold=True, color=ACCENT_BLUE)
add_text_box(slide,
    "Provera demonstrates that hybrid ML + agentic AI systems can significantly improve Medicare fraud detection. "
    "By combining graph ML for structural pattern recognition with LLM agents for investigation and explanation, "
    "we achieve both high accuracy (0.91 ROC-AUC) and full explainability (100% factual accuracy).",
    0.8, 6.2, 11.733, 0.6, font_size=11, color=GRAY)


# ============ Save ============
output_path = "/Users/nikitaravi/Desktop/IntelliMed/medigraph/presentation/Provera_Capstone.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: 14 (Title + 13 content)")
